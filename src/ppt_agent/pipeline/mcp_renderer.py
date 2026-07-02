"""Render presentations via the Office-PowerPoint MCP Server.

DesignDNA-driven renderer: all geometry, colors, and fonts come from
the DesignDNA profile extracted from a reference template.

Split across three collaborators:
  - mcp_client.MCPClient  — stdio transport + tool-error accounting
  - geometry              — pure DNA→layout math and color helpers
  - layouts               — one builder per SlideLayout (DISPATCH table)
This module keeps the public MCPRenderer API and the drawing primitives
the layout builders compose.
"""

from __future__ import annotations

import asyncio
from pathlib import Path

from mcp import StdioServerParameters
from rich.console import Console

from ppt_agent.models import DesignDNA, Presentation, SlideContent, SlideLayout
from ppt_agent.pipeline import layouts
from ppt_agent.pipeline.geometry import (
    Geometry,
    aspect_fit,
    fit_font_size,
    hex_to_rgb as _hex,  # re-exported; external callers/tests import from here
)
from ppt_agent.pipeline.mcp_client import MCPClient
from ppt_agent.utils.fonts import font_for_text, get_font_pair

console = Console()


class MCPRenderer:
    """DesignDNA-driven presentation renderer using the MCP Server tool palette."""

    _RATIOS = {
        "16:9": (13.333, 7.5),
        "4:3": (10.0, 7.5),
    }

    SPACING = 0.12
    CARD_PAD = 0.15
    ELEM_GAP = 0.2

    _TRANSITIONS = {
        SlideLayout.TITLE: {"transition_type": "fade", "duration": 1.0},
        SlideLayout.SECTION: {"transition_type": "fade", "duration": 0.8},
        SlideLayout.CLOSING: {"transition_type": "fade", "duration": 1.0},
    }
    _DEFAULT_TRANSITION = {"transition_type": "fade", "duration": 0.5}

    def __init__(self, design_dna: DesignDNA | None = None, aspect_ratio: str = "16:9"):
        self.dna = design_dna or DesignDNA()
        self._server_params = StdioServerParameters(
            command="uvx",
            args=["--from", "office-powerpoint-mcp-server", "ppt_mcp_server"],
        )
        self._section_num = 0
        self._font_title, self._font_body = get_font_pair()
        self._client: MCPClient | None = None

        dims = self._RATIOS.get(aspect_ratio, self._RATIOS["16:9"])
        self.W, self.H = dims
        self._compute_geometry()

    def _compute_geometry(self) -> None:
        """Derive all layout positions from DesignDNA ratios."""
        self.geo = Geometry(self.dna, self.W, self.H)
        for name in Geometry.FIELDS:
            setattr(self, name, getattr(self.geo, name))

    def render(self, presentation: Presentation, output_path: str) -> str:
        path = Path(output_path)
        path.parent.mkdir(parents=True, exist_ok=True)
        asyncio.run(self._render_async(presentation, str(path.resolve())))
        self._apply_post_process(path, presentation)
        self._report_degradation()
        return str(path)

    def _report_degradation(self) -> None:
        """Surface per-shape failures that were degraded rather than fatal."""
        if not self._client:
            return
        stats = self._client.stats
        if stats["tool_errors"] or stats["degraded_shapes"]:
            console.print(
                f"  [yellow]渲染降级：{stats['tool_errors']} 次工具调用报错，"
                f"{stats['degraded_shapes']} 个元素被跳过（转场/阴影/图片等）。"
                f"如成品明显异常，请检查 office-powerpoint-mcp-server 版本[/yellow]"
            )

    def _apply_post_process(self, path: Path, presentation: Presentation) -> None:
        """Fix up the saved deck in a single python-pptx pass.

        1. Set page size to match render geometry. The MCP server creates 4:3
           decks and exposes no sizing tool, while all geometry here targets
           self.W x self.H; shape positions are absolute EMU, so resizing the
           page after save is lossless.
        2. Inject speaker notes. The MCP server's manage_text "notes" target
           degrades to a stray on-canvas textbox, so notes are written here
           straight into each slide's notes placeholder instead.
        3. Inject fade transitions. The MCP server's manage_slide_transitions
           tool is a placeholder no-op, so the transition XML is written here.
        """
        from pptx import Presentation as PptxPresentation
        from pptx.util import Inches

        pptx = PptxPresentation(str(path))
        pptx.slide_width = Inches(self.W)
        pptx.slide_height = Inches(self.H)

        # Slides are built sequentially, so pptx order matches presentation order.
        for slide, sc in zip(pptx.slides, presentation.slides):
            if sc.notes:
                slide.notes_slide.notes_text_frame.text = sc.notes
            trans = self._TRANSITIONS.get(sc.layout, self._DEFAULT_TRANSITION)
            self._inject_transition(slide, trans["duration"])

        pptx.save(str(path))

    @staticmethod
    def _inject_transition(slide, duration: float) -> None:
        """Write a <p:transition><p:fade/></p:transition> into the slide XML."""
        from pptx.oxml.ns import qn

        sld = slide.element
        for old in sld.findall(qn("p:transition")):
            sld.remove(old)
        trans = sld.makeelement(qn("p:transition"), {})
        # PowerPoint speed buckets: fast≈0.5s, med≈0.75s, slow≈1s.
        trans.set("spd", "fast" if duration <= 0.5 else "med" if duration <= 0.8 else "slow")
        trans.append(sld.makeelement(qn("p:fade"), {}))
        sld.append(trans)  # schema order: after cSld/clrMapOvr; no timing element exists

    # ── async core ─────────────────────────────────────────────

    async def _render_async(self, pres: Presentation, out: str) -> None:
        async with MCPClient(self._server_params) as client:
            self._client = client
            await self._build(client, pres, out)

    async def _build(self, s: MCPClient, pres: Presentation, out: str) -> None:
        r = await s.call("create_presentation", {"id": "ppt_agent"})
        s.presentation_id = r.get("presentation_id", "ppt_agent")
        self._section_num = 0

        await s.call("set_core_properties", {
            "title": pres.outline.title,
            "author": pres.author or "",
        })

        total = len(pres.slides)

        # Each MCP session starts from an empty presentation, so a failed render
        # cannot be resumed mid-deck; every run rebuilds all slides from scratch.
        for i, sc in enumerate(pres.slides):
            try:
                idx = await self._dispatch(s, sc, pres, i)
                if sc.layout not in (SlideLayout.TITLE, SlideLayout.CLOSING, SlideLayout.SECTION):
                    await self._footer(s, idx, i + 1, total)
            except Exception as e:
                await s.call("save_presentation", {"file_path": out})
                raise RuntimeError(
                    f"Render failed at slide {i + 1}/{total} ({sc.layout.value}): {e}. "
                    f"Partial deck (slides 1-{i}) saved for inspection; "
                    f"re-running renders the full deck from scratch."
                ) from e

        await s.call("save_presentation", {"file_path": out})

    async def _dispatch(self, s: MCPClient, sc: SlideContent, pres: Presentation, i: int) -> int:
        fn = layouts.DISPATCH.get(sc.layout, layouts.slide_image_focus)
        idx = await fn(self, s, sc, pres)

        # Speaker notes and slide transitions are injected post-render via
        # python-pptx (_apply_post_process): the MCP server's manage_text
        # "notes" target degrades to a stray on-canvas textbox, and its
        # manage_slide_transitions is a documented placeholder no-op.

        return idx

    # ── DNA color helpers ─────────────────────────────────────

    @property
    def _primary(self) -> list[int]:
        return _hex(self.dna.primary_color)

    @property
    def _accent(self) -> list[int]:
        return _hex(self.dna.accent_color)

    @property
    def _bg2(self) -> list[int]:
        return _hex(self.dna.secondary_color)

    @property
    def _text_on_primary(self) -> list[int]:
        return _hex(self.dna.text_on_primary)

    @property
    def _text_on_light(self) -> list[int]:
        return _hex(self.dna.text_on_light)

    @property
    def _text_muted(self) -> list[int]:
        return _hex(self.dna.text_muted)

    @property
    def _frame_color(self) -> list[int]:
        return _hex(self.dna.frame_border_color)

    def _sz(self, tier: str) -> int:
        return self.dna.size_hierarchy.get(tier, 20)

    # ── visual primitives ──────────────────────────────────────

    async def _blank(self, s: MCPClient) -> int:
        r = await s.call("add_slide", {"layout_index": 6})
        return r.get("slide_index", 0)

    async def _rect(self, s, idx, left, top, w, h, fill, **kw) -> dict:
        kw.setdefault("line_color", fill)
        args = {
            "slide_index": idx, "shape_type": "rectangle",
            "left": left, "top": top, "width": w, "height": h,
            "fill_color": fill,
        }
        args.update(kw)
        return await s.call("add_shape", args)

    async def _rrect(self, s, idx, left, top, w, h, fill, **kw) -> dict:
        kw.setdefault("line_color", fill)
        args = {
            "slide_index": idx, "shape_type": "rounded_rectangle",
            "left": left, "top": top, "width": w, "height": h,
            "fill_color": fill,
        }
        args.update(kw)
        return await s.call("add_shape", args)

    async def _text(self, s, idx, left, top, w, h, text, size=18, color=None,
                    bold=False, align="left", font=None) -> dict:
        color = color or self._text_on_light
        if font is None:
            font = font_for_text(text)
        size = fit_font_size(text, w, h, size)
        return await s.call("manage_text", {
            "slide_index": idx, "operation": "add",
            "left": left, "top": top, "width": w, "height": h,
            "text": text, "font_size": size, "bold": bold,
            "font_name": font, "color": color, "alignment": align,
        })

    async def _shadow(self, s, idx, shape_idx) -> None:
        try:
            await s.call("apply_picture_effects", {
                "slide_index": idx, "shape_index": shape_idx,
                "effects": {
                    "shadow": {
                        "shadow_type": "outer", "blur_radius": 6.0,
                        "distance": 3.0, "direction": 315.0,
                        "color": [0, 0, 0], "transparency": 0.7,
                    }
                },
            })
        except Exception:
            s.stats["degraded_shapes"] += 1

    # ── composite primitives ───────────────────────────────────

    async def _header_band(self, s, idx, title) -> None:
        size = self._sz("header")
        await self._rect(s, idx, 0, self.HDR_TOP, self.W, self.HDR_H, self._primary)
        await self._text(s, idx, self.M_LEFT, self.HDR_TOP + 0.08,
                         self.CW, self.HDR_H - 0.16,
                         title, size=size, bold=True,
                         color=self._text_on_primary, align="center")

    async def _key_statement(self, s, idx, text) -> None:
        if not text:
            return
        size = self._sz("key_statement")
        await self._text(s, idx, 0, self.KEY_Y, self.W, self.KEY_H,
                         text, size=size, bold=True,
                         color=self._text_on_light, align="center")

    async def _frame(self, s, idx, left, top, w, h, line_color=None) -> dict:
        lc = line_color or self._frame_color
        return await self._rrect(s, idx, left, top, w, h, [255, 255, 255],
                                 line_color=lc, line_width=self.dna.frame_border_width)

    async def _annotation_list(self, s, idx, items, left, top, w,
                               size=None, spacing=None, height=None) -> float:
        """Render a bullet list; with `height`, distribute and center vertically."""
        size = size or self._sz("body")
        items = items[:5]
        n = max(len(items), 1)
        if height:
            sp = spacing or min(1.0, height / n)
            y = top + max(0.0, (height - n * sp) / 2)
        else:
            sp = spacing or min(0.7, self.BODY_H / n)
            y = top
        for item in items:
            await self._text(s, idx, left, y, w, min(0.9, sp),
                             f"•  {item}", size=size, bold=False,
                             color=self._text_on_light)
            y += sp
        return y

    async def _citation(self, s, idx, text) -> None:
        if not text:
            return
        await self._text(s, idx, self.M_LEFT, self.CITE_Y, self.CW, 0.35,
                         text, size=self._sz("citation"), bold=False,
                         color=self._text_muted)

    async def _footer(self, s, idx, page, total) -> None:
        await self._text(s, idx, self.W - 1.4, self.H - 0.35, 1.0, 0.25,
                         f"{page} / {total}", size=10, bold=False,
                         color=self._text_muted, align="right")

    def _image_dir(self) -> str:
        """Directory for renderer-generated assets (matplotlib charts, etc.)."""
        import tempfile

        d = Path(tempfile.gettempdir()) / "ppt_agent_images"
        d.mkdir(parents=True, exist_ok=True)
        return str(d)

    async def _insert_image(self, s, idx, img_spec, left, top, w, h) -> None:
        """Insert an image file or a placeholder frame."""
        if img_spec and img_spec.local_path and Path(img_spec.local_path).exists():
            try:
                fl, ft, fw, fh = aspect_fit(img_spec.local_path, left, top, w, h)
                r = await s.call("manage_image", {
                    "slide_index": idx, "operation": "add",
                    "image_source": img_spec.local_path, "source_type": "file",
                    "left": fl, "top": ft, "width": fw, "height": fh,
                })
                if "error" in r:
                    raise RuntimeError(r["error"])
                if img_spec.caption:
                    await self._text(s, idx, left, top + h + 0.05, w, 0.35,
                                     img_spec.caption, size=self._sz("caption"),
                                     bold=False, color=self._text_muted, align="center")
                return
            except Exception:
                s.stats["degraded_shapes"] += 1

        await self._frame(s, idx, left, top, w, h)
        desc = img_spec.description if img_spec else "Image"
        await self._text(s, idx, left + 0.3, top + h / 2 - 0.3, w - 0.6, 0.6,
                         f"[ {desc} ]", size=self._sz("caption"), bold=True,
                         color=self._text_muted, align="center")
