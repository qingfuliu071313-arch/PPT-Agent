"""Extract DesignDNA from an existing PPTX file using python-pptx."""

from __future__ import annotations

from collections import Counter
from statistics import median

from pptx import Presentation as PptxPresentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from ppt_agent.models import DesignDNA
from ppt_agent.template_learning.spatial_analyzer import (
    cluster_font_sizes,
    compute_image_ratios,
    detect_content_margins,
    detect_header_bar,
    find_dark_text_color,
    find_text_color_on_dark,
    most_common_non_neutral,
)


class TemplateExtractor:
    """Extract a complete DesignDNA from a PPTX template."""

    def extract(self, pptx_path: str) -> DesignDNA:
        prs = PptxPresentation(pptx_path)
        slide_w = prs.slide_width
        slide_h = prs.slide_height

        colors: list[tuple[str, str]] = []
        fonts: list[tuple[str, float, bool]] = []
        shapes: list[dict] = []
        images: list[dict] = []
        alignments: list[str] = []
        line_spacings: list[float] = []
        slides_with_images = 0

        for slide_idx, slide in enumerate(prs.slides):
            slide_has_image = False

            for shape in slide.shapes:
                shape_info = {
                    "slide_index": slide_idx,
                    "type": self._shape_type_name(shape),
                    "left": shape.left or 0,
                    "top": shape.top or 0,
                    "width": shape.width or 0,
                    "height": shape.height or 0,
                    "has_text": shape.has_text_frame,
                }

                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    slide_has_image = True
                    images.append(shape_info)

                if shape.has_text_frame:
                    char_count = 0
                    max_size = 0
                    for para in shape.text_frame.paragraphs:
                        if para.alignment is not None:
                            alignments.append(str(para.alignment))
                        try:
                            sp = para.line_spacing
                            if sp is not None and isinstance(sp, (int, float)):
                                line_spacings.append(float(sp))
                        except (AttributeError, TypeError):
                            pass

                        for run in para.runs:
                            char_count += len(run.text)
                            if run.font.name:
                                size = run.font.size.pt if run.font.size else 0
                                fonts.append((run.font.name, size, bool(run.font.bold)))
                                if size > max_size:
                                    max_size = size
                            self._collect_color(run, colors)

                    shape_info["char_count"] = char_count
                    shape_info["font_size"] = max_size

                self._collect_shape_color(shape, colors)
                shapes.append(shape_info)

            if slide_has_image:
                slides_with_images += 1

        total_slides = len(list(prs.slides))

        fill_colors = [c for c, ctx in colors if ctx == "fill"]
        primary = most_common_non_neutral(fill_colors)
        text_on_primary = find_text_color_on_dark(colors)
        text_on_light = find_dark_text_color(colors)

        font_names = Counter(name for name, _, _ in fonts if name)
        primary_font = font_names.most_common(1)[0][0] if font_names else "Microsoft YaHei"

        all_sizes = [s for _, s, _ in fonts if s > 0]
        size_hierarchy = cluster_font_sizes(all_sizes)

        bold_count = sum(1 for _, _, b in fonts if b)
        bold_ratio = bold_count / len(fonts) if fonts else 0.5

        align_counter = Counter(alignments)
        dominant_alignment = "center"
        if align_counter:
            top_align = align_counter.most_common(1)[0][0]
            mapping = {"CENTER (2)": "center", "LEFT (0)": "left", "JUSTIFY (3)": "justify"}
            for key, val in mapping.items():
                if key in top_align:
                    dominant_alignment = val
                    break

        avg_spacing = median(line_spacings) if line_spacings else 1.5

        header_bar = detect_header_bar(shapes, slide_w, slide_h)
        content_margins = detect_content_margins(shapes, header_bar, slide_w, slide_h)
        image_area_ratios = compute_image_ratios(images, shapes, slide_w, slide_h)

        secondary = self._find_secondary_color(fill_colors, primary)

        frame_colors = [c for c, ctx in colors if ctx == "line"]
        frame_border_color = most_common_non_neutral(frame_colors, primary)

        return DesignDNA(
            name="extracted",
            primary_color=primary,
            secondary_color=secondary,
            accent_color=self._find_accent(fill_colors, primary, secondary),
            text_on_primary=text_on_primary,
            text_on_light=text_on_light,
            text_muted=self._find_muted_color(colors),
            font_primary=primary_font,
            font_secondary=font_names.most_common(2)[-1][0] if len(font_names) >= 2 else primary_font,
            bold_ratio=round(bold_ratio, 2),
            size_hierarchy=size_hierarchy,
            dominant_alignment=dominant_alignment,
            line_spacing=round(avg_spacing, 1),
            header_bar=header_bar,
            content_margins=content_margins,
            image_area_ratios=image_area_ratios,
            frame_border_color=frame_border_color,
            frame_border_width=1.0,
            avg_images_per_slide=round(len(images) / max(total_slides, 1), 1),
            image_slide_percentage=round(slides_with_images / max(total_slides, 1), 2),
        )

    def _shape_type_name(self, shape) -> str:
        try:
            return shape.shape_type.name if hasattr(shape.shape_type, "name") else str(shape.shape_type)
        except Exception:
            return "UNKNOWN"

    def _collect_color(self, run, colors: list[tuple[str, str]]) -> None:
        try:
            if run.font.color and run.font.color.type is not None:
                rgb = run.font.color.rgb
                if rgb is not None:
                    colors.append((str(rgb), "text"))
        except (AttributeError, TypeError):
            pass

    def _collect_shape_color(self, shape, colors: list[tuple[str, str]]) -> None:
        try:
            if hasattr(shape, "fill") and shape.fill and shape.fill.type is not None:
                rgb = shape.fill.fore_color.rgb
                if rgb is not None:
                    colors.append((str(rgb), "fill"))
        except (AttributeError, TypeError):
            pass
        try:
            if shape.line and shape.line.fill and shape.line.fill.type is not None:
                rgb = shape.line.color.rgb
                if rgb is not None:
                    colors.append((str(rgb), "line"))
        except (AttributeError, TypeError):
            pass

    def _find_secondary_color(self, fill_colors: list[str], primary: str) -> str:
        filtered = [c for c in fill_colors if c != primary]
        light = []
        for c in filtered:
            try:
                r, g, b = int(c[:2], 16), int(c[2:4], 16), int(c[4:6], 16)
                if r + g + b > 600:
                    light.append(c)
            except (ValueError, IndexError):
                continue
        if light:
            return Counter(light).most_common(1)[0][0]
        return "F5F5F5"

    def _find_accent(self, fill_colors: list[str], primary: str, secondary: str) -> str:
        filtered = [c for c in fill_colors if c not in (primary, secondary, "FFFFFF", "000000")]
        if filtered:
            return Counter(filtered).most_common(1)[0][0]
        return primary

    def _find_muted_color(self, colors: list[tuple[str, str]]) -> str:
        text_colors = [c for c, ctx in colors if ctx == "text"]
        muted = []
        for c in text_colors:
            try:
                r, g, b = int(c[:2], 16), int(c[2:4], 16), int(c[4:6], 16)
                if 200 < r + g + b < 500:
                    muted.append(c)
            except (ValueError, IndexError):
                continue
        if muted:
            return Counter(muted).most_common(1)[0][0]
        return "666666"
