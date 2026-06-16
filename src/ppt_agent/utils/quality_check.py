"""Post-render quality verification for generated PPTX files."""

from __future__ import annotations

from pathlib import Path
from pptx import Presentation as PptxPresentation


class QualityIssue:
    def __init__(self, slide: int, severity: str, message: str):
        self.slide = slide
        self.severity = severity  # "error", "warning", "info"
        self.message = message

    def __str__(self) -> str:
        return f"[{self.severity.upper()}] Slide {self.slide}: {self.message}"


def verify_pptx(pptx_path: str, expected_slides: int = 0) -> list[QualityIssue]:
    """Verify a generated PPTX file for completeness and correctness."""
    path = Path(pptx_path)
    issues: list[QualityIssue] = []

    if not path.exists():
        return [QualityIssue(0, "error", f"File not found: {pptx_path}")]

    if path.stat().st_size < 5000:
        issues.append(QualityIssue(0, "warning", f"File very small ({path.stat().st_size} bytes)"))

    try:
        prs = PptxPresentation(str(path))
    except Exception as e:
        return [QualityIssue(0, "error", f"Cannot open file: {e}")]

    actual_slides = len(prs.slides)

    if actual_slides == 0:
        issues.append(QualityIssue(0, "error", "Presentation has no slides"))
        return issues

    if expected_slides > 0 and actual_slides != expected_slides:
        issues.append(QualityIssue(0, "warning",
                                    f"Expected {expected_slides} slides, got {actual_slides}"))

    for i, slide in enumerate(prs.slides, 1):
        shape_count = len(slide.shapes)

        if shape_count == 0:
            issues.append(QualityIssue(i, "error", "Empty slide (no shapes)"))
            continue

        if shape_count < 3:
            issues.append(QualityIssue(i, "warning", f"Very few shapes ({shape_count})"))

        has_text = False
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    has_text = True
                    if len(text) > 500:
                        issues.append(QualityIssue(i, "warning",
                                                    f"Very long text block ({len(text)} chars)"))

        if not has_text:
            issues.append(QualityIssue(i, "warning", "No visible text on slide"))

        has_notes = (slide.has_notes_slide and
                     slide.notes_slide.notes_text_frame.text.strip())
        if not has_notes and i not in (1, actual_slides):
            issues.append(QualityIssue(i, "info", "Missing speaker notes"))

    return issues


def quality_report(pptx_path: str, expected_slides: int = 0) -> str:
    """Generate a human-readable quality report."""
    issues = verify_pptx(pptx_path, expected_slides)

    if not issues:
        return "✓ Quality check passed — no issues found."

    errors = [i for i in issues if i.severity == "error"]
    warnings = [i for i in issues if i.severity == "warning"]
    infos = [i for i in issues if i.severity == "info"]

    lines = [f"Quality Report: {len(issues)} issue(s) found"]
    if errors:
        lines.append(f"  Errors: {len(errors)}")
    if warnings:
        lines.append(f"  Warnings: {len(warnings)}")
    if infos:
        lines.append(f"  Info: {len(infos)}")
    lines.append("")
    for i in issues:
        lines.append(f"  {i}")

    return "\n".join(lines)
