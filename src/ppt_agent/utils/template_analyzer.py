"""Analyze existing PPT files to extract style information.

Uses MCP Server tools to inspect slide masters, colors, fonts, and layouts.
"""

from __future__ import annotations

import asyncio
import json

from mcp import ClientSession, StdioServerParameters
from mcp.client.stdio import stdio_client

from ppt_agent.models import StyleConfig, ThemeConfig


class TemplateAnalyzer:
    """Extract style information from an existing PPTX file."""

    def __init__(self):
        self._server_params = StdioServerParameters(
            command="uvx",
            args=["--from", "office-powerpoint-mcp-server", "ppt_mcp_server"],
        )
        self._pres_id: str = ""

    def analyze(self, pptx_path: str) -> dict:
        """Analyze a PPTX file and return comprehensive style information."""
        return asyncio.run(self._analyze_async(pptx_path))

    async def _analyze_async(self, pptx_path: str) -> dict:
        async with stdio_client(self._server_params) as streams:
            async with ClientSession(*streams) as s:
                await s.initialize()

                r = await self._call(s, "open_presentation", {"file_path": pptx_path})
                self._pres_id = r.get("presentation_id", "")

                info = await self._call(s, "get_presentation_info", {})
                masters = await self._call(s, "manage_slide_masters", {"operation": "list"})

                slides_info = []
                slide_count = info.get("slide_count", 0)
                for i in range(min(slide_count, 5)):
                    si = await self._call(s, "get_slide_info", {"slide_index": i})
                    slides_info.append(si)

                return {
                    "presentation": info,
                    "slide_masters": masters,
                    "sample_slides": slides_info,
                    "slide_count": slide_count,
                }

    async def _call(self, s: ClientSession, tool: str, args: dict) -> dict:
        if self._pres_id and "presentation_id" not in args:
            args = {**args, "presentation_id": self._pres_id}
        try:
            result = await s.call_tool(tool, arguments=args)
            if result.content:
                for block in result.content:
                    if hasattr(block, "text"):
                        try:
                            return json.loads(block.text)
                        except (json.JSONDecodeError, AttributeError):
                            return {"raw": block.text}
        except Exception:
            return {}
        return {}

    def extract_style(self, pptx_path: str) -> StyleConfig:
        """Extract a StyleConfig from an existing PPTX file.

        Falls back to python-pptx direct analysis for color/font extraction.
        """
        try:
            return self._extract_via_pptx(pptx_path)
        except Exception:
            return StyleConfig()

    def _extract_via_pptx(self, pptx_path: str) -> StyleConfig:
        """Use python-pptx directly for reliable style extraction."""
        from pptx import Presentation as PptxPresentation
        from pptx.util import Pt

        prs = PptxPresentation(pptx_path)

        colors = []
        fonts = set()
        font_sizes = {"title": [], "body": []}

        for slide in list(prs.slides)[:5]:
            bg = slide.background
            if bg.fill and bg.fill.type is not None:
                try:
                    c = bg.fill.fore_color.rgb
                    colors.append(str(c))
                except Exception:
                    pass

            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if run.font.name:
                                fonts.add(run.font.name)
                            try:
                                if run.font.color and run.font.color.type is not None:
                                    rgb = run.font.color.rgb
                                    if rgb is not None:
                                        colors.append(str(rgb))
                            except (AttributeError, TypeError):
                                pass
                            if run.font.size:
                                pts = run.font.size.pt
                                if pts >= 24:
                                    font_sizes["title"].append(pts)
                                else:
                                    font_sizes["body"].append(pts)

                try:
                    if hasattr(shape, "fill") and shape.fill and shape.fill.type is not None:
                        c = shape.fill.fore_color.rgb
                        if c is not None:
                            colors.append(str(c))
                except (AttributeError, TypeError):
                    pass

        primary = _most_common(colors, "1F4E79")
        text_color = _find_dark_color(colors, "333333")
        font_list = sorted(fonts)

        title_size = int(sum(font_sizes["title"]) / len(font_sizes["title"])) if font_sizes["title"] else 28
        body_size = int(sum(font_sizes["body"]) / len(font_sizes["body"])) if font_sizes["body"] else 18

        return StyleConfig(
            primary_color=primary,
            accent_color=_secondary_color(colors, primary),
            background_color="FFFFFF",
            text_color=text_color,
            font_title=font_list[0] if font_list else "PingFang SC",
            font_body=font_list[0] if font_list else "PingFang SC",
            font_size_title=title_size,
            font_size_body=body_size,
        )

    def extract_theme(self, pptx_path: str, name: str = "custom") -> ThemeConfig:
        """Extract a ThemeConfig suitable for the MCPRenderer."""
        style = self.extract_style(pptx_path)
        return ThemeConfig(
            name=name,
            primary_color=style.primary_color,
            secondary_color="F2F4F7",
            accent_color=style.accent_color,
            text_dark=style.text_color,
            text_light="FFFFFF",
            text_muted="718096",
            font_title=style.font_title,
            font_body=style.font_body,
        )


def _most_common(colors: list[str], default: str) -> str:
    if not colors:
        return default
    filtered = [c for c in colors if c not in ("FFFFFF", "000000", "ffffff", "000000")]
    if not filtered:
        return default
    from collections import Counter
    return Counter(filtered).most_common(1)[0][0]


def _find_dark_color(colors: list[str], default: str) -> str:
    for c in colors:
        try:
            r, g, b = int(c[:2], 16), int(c[2:4], 16), int(c[4:6], 16)
            if r + g + b < 300:
                return c
        except (ValueError, IndexError):
            continue
    return default


def _secondary_color(colors: list[str], primary: str) -> str:
    filtered = [c for c in colors if c not in ("FFFFFF", "000000", primary)]
    if not filtered:
        return primary
    from collections import Counter
    return Counter(filtered).most_common(1)[0][0]
